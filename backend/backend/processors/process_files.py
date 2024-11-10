import os
import shutil
import subprocess
import tempfile
from datetime import datetime
from typing import List

import fitz
from fastapi import UploadFile
from llama_index.core import Document, SimpleDirectoryReader
from pptx import Presentation

from backend.config import logger

from .utils import (
    async_timeout,
    describe_image,
    extract_text_around_item,
    is_graph,
    process_graph,
    process_text_blocks,
    save_uploaded_file,
)

# TODO: Split into separate files


async def get_pdf_documents(
    pdf_file: UploadFile, research_id: str, reference_id: str
) -> List[Document]:
    """Process a PDF file and extract text, tables, and images."""
    all_pdf_documents = []
    ongoing_tables = {}

    try:
        logger.info(f"Reading PDF file: {pdf_file.filename}")
        pdf_content = await pdf_file.read()  # Await the read coroutine
        f = fitz.open(stream=pdf_content, filetype="pdf")
        logger.info(f"Successfully opened PDF file: {pdf_file.filename}")
    except Exception as e:
        logger.error(f"Error opening or processing the PDF file: {e}")
        return []

    for i in range(len(f)):
        page = f[i]
        logger.info(f"Processing page {i+1} of {len(f)}")
        text_blocks = [
            block
            for block in page.get_text("blocks", sort=True)
            if block[-1] == 0
            and not (
                block[1] < page.rect.height * 0.1 or block[3] > page.rect.height * 0.9
            )
        ]
        grouped_text_blocks = process_text_blocks(text_blocks)

        table_docs, table_bboxes, ongoing_tables = await parse_all_tables(
            pdf_file.filename, page, i, text_blocks, ongoing_tables
        )
        all_pdf_documents.extend(table_docs)
        logger.info(f"Extracted {len(table_docs)} tables from page {i+1}")

        image_docs = await parse_all_images(pdf_file.filename, page, i, text_blocks)
        all_pdf_documents.extend(image_docs)
        logger.info(f"Extracted {len(image_docs)} images from page {i+1}")

        for text_block_ctr, (heading_block, content) in enumerate(
            grouped_text_blocks, 1
        ):
            heading_bbox = fitz.Rect(heading_block[:4])
            if not any(
                heading_bbox.intersects(table_bbox) for table_bbox in table_bboxes
            ):
                bbox = {
                    "x1": heading_block[0],
                    "y1": heading_block[1],
                    "x2": heading_block[2],
                    "x3": heading_block[3],
                }
                text_doc = Document(
                    text=f"{heading_block[4]}\n{content}",
                    metadata={
                        **bbox,
                        "type": "text",
                        "page_num": i,
                        "source": f"{pdf_file.filename[:-4]}-page{i}-block{text_block_ctr}",
                        "research_id": research_id,
                        "reference_id": reference_id,
                        "filename": pdf_file.filename,
                        "created_at": datetime.utcnow().isoformat(),
                    },
                )
                all_pdf_documents.append(text_doc)
        logger.info(f"Extracted {len(grouped_text_blocks)} text blocks from page {i+1}")

    f.close()
    logger.info(f"Completed processing PDF file: {pdf_file.filename}")
    return all_pdf_documents


@async_timeout(20)
async def process_single_table(tab, text_blocks, page, pagenum, filename, table_docs):
    """Process a single table with timeout"""
    try:
        if not tab.header.external:
            pandas_df = tab.to_pandas()
            # ...existing table processing code...
            tablerefdir = os.path.join(os.getcwd(), "vectorstore/table_references")
            os.makedirs(tablerefdir, exist_ok=True)

            df_xlsx_path = os.path.join(
                tablerefdir, f"table{len(table_docs)+1}-page{pagenum}.xlsx"
            )
            pandas_df.to_excel(df_xlsx_path)
            bbox = fitz.Rect(tab.bbox)

            before_text, after_text = extract_text_around_item(
                text_blocks, bbox, page.rect.height
            )
            table_img = page.get_pixmap(clip=bbox)
            table_img_path = os.path.join(
                tablerefdir, f"table{len(table_docs)+1}-page{pagenum}.jpg"
            )
            table_img.save(table_img_path)
            description = process_graph(table_img.tobytes())

            caption = (
                before_text.replace("\n", " ")
                + description
                + after_text.replace("\n", " ")
            )
            if before_text == "" and after_text == "":
                caption = " ".join(tab.header.names)

            table_metadata = {
                "source": f"{filename[:-4]}-page{pagenum}-table{len(table_docs)+1}",
                "dataframe": df_xlsx_path,
                "image": table_img_path,
                "caption": caption,
                "type": "table",
                "page_num": pagenum,
            }

            all_cols = ", ".join(list(pandas_df.columns.values))
            return Document(
                text=f"This is a table with the caption: {caption}\nThe columns are {all_cols}",
                metadata=table_metadata,
            )
    except Exception as e:
        logger.error(f"Error processing table: {e}")
        return None


async def parse_all_tables(filename, page, pagenum, text_blocks, ongoing_tables):
    """Extract tables from a PDF page."""
    table_docs = []
    table_bboxes = []

    try:
        logger.info(f"Finding tables on page {pagenum+1}")
        tables = page.find_tables(
            horizontal_strategy="lines_strict", vertical_strategy="lines_strict"
        )
        logger.info(f"Found {len(tables.tables)} tables on page {pagenum+1}")

        for tab in tables.tables:
            doc = await process_single_table(
                tab, text_blocks, page, pagenum, filename, table_docs
            )
            if doc:
                table_docs.append(doc)
                table_bboxes.append(fitz.Rect(tab.bbox))

    except Exception as e:
        logger.error(f"Error during table extraction: {e}")

    return table_docs, table_bboxes, ongoing_tables


@async_timeout(20)
async def process_single_image(image_info, page, pagenum, text_blocks, filename):
    """Process a single image with timeout"""
    xref = image_info["xref"]
    if xref == 0:
        return None

    # ...existing image processing code...
    img_bbox = fitz.Rect(image_info["bbox"])
    if img_bbox.width < page.rect.width / 20 or img_bbox.height < page.rect.height / 20:
        return None

    try:
        extracted_image = page.parent.extract_image(xref)
        image_data = extracted_image["image"]

        # Process image and create metadata
        imgrefpath = os.path.join(os.getcwd(), "vectorstore/image_references")
        os.makedirs(imgrefpath, exist_ok=True)
        image_path = os.path.join(imgrefpath, f"image{xref}-page{pagenum}.png")

        with open(image_path, "wb") as img_file:
            img_file.write(image_data)

        before_text, after_text = extract_text_around_item(
            text_blocks, img_bbox, page.rect.height
        )
        if before_text == "" and after_text == "":
            return None

        image_description = " "
        if is_graph(image_data):
            image_description = process_graph(image_data)

        caption = (
            before_text.replace("\n", " ")
            + image_description
            + after_text.replace("\n", " ")
        )

        image_metadata = {
            "source": f"{filename[:-4]}-page{pagenum}-image{xref}",
            "image": image_path,
            "caption": caption,
            "type": "image",
            "page_num": pagenum,
        }

        return Document(
            text=f"This is an image with the caption: " + caption,
            metadata=image_metadata,
        )
    except Exception as e:
        logger.error(f"Error processing image: {e}")
        return None


async def parse_all_images(filename, page, pagenum, text_blocks):
    """Extract images from a PDF page."""
    image_docs = []
    image_info_list = page.get_image_info(xrefs=True)

    logger.info(f"Finding images on page {pagenum+1}")
    for image_info in image_info_list:
        doc = await process_single_image(
            image_info, page, pagenum, text_blocks, filename
        )
        if doc:
            image_docs.append(doc)

    logger.info(f"Found {len(image_docs)} images on page {pagenum+1}")
    return image_docs


def process_ppt_file(
    ppt_path: str, research_id: str, reference_id: str, filename: str
) -> List[Document]:
    """Process a PowerPoint file."""
    pdf_path = convert_ppt_to_pdf(ppt_path)
    images_data = convert_pdf_to_images(pdf_path)
    slide_texts = extract_text_and_notes_from_ppt(ppt_path)
    processed_data = []

    for (image_path, page_num), (slide_text, notes) in zip(images_data, slide_texts):
        if notes:
            notes = "\n\nThe speaker notes for this slide are: " + notes

        with open(image_path, "rb") as image_file:
            image_content = image_file.read()

        image_description = " "
        if is_graph(image_content):
            image_description = process_graph(image_content)

        image_metadata = {
            "source": f"{os.path.basename(ppt_path)}",
            "image": image_path,
            "caption": slide_text + image_description + notes,
            "type": "image",
            "page_num": page_num,
            "research_id": research_id,
            "reference_id": reference_id,
            "filename": filename,
            "created_at": datetime.utcnow().isoformat(),
        }
        processed_data.append(
            Document(
                text="This is a slide with the text: " + slide_text + image_description,
                metadata=image_metadata,
            )
        )

    return processed_data


def convert_ppt_to_pdf(ppt_path):
    """Convert a PowerPoint file to PDF using LibreOffice."""
    base_name = os.path.basename(ppt_path)
    ppt_name_without_ext = os.path.splitext(base_name)[0].replace(" ", "_")
    new_dir_path = os.path.abspath("vectorstore/ppt_references")
    os.makedirs(new_dir_path, exist_ok=True)
    pdf_path = os.path.join(new_dir_path, f"{ppt_name_without_ext}.pdf")
    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        new_dir_path,
        ppt_path,
    ]
    subprocess.run(command, check=True)
    return pdf_path


def convert_pdf_to_images(pdf_path):
    """Convert a PDF file to a series of images using PyMuPDF."""
    doc = fitz.open(pdf_path)
    base_name = os.path.basename(pdf_path)
    pdf_name_without_ext = os.path.splitext(base_name)[0].replace(" ", "_")
    new_dir_path = os.path.join(os.getcwd(), "vectorstore/ppt_references")
    os.makedirs(new_dir_path, exist_ok=True)
    image_paths = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap()
        output_image_path = os.path.join(
            new_dir_path, f"{pdf_name_without_ext}_{page_num:04d}.png"
        )
        pix.save(output_image_path)
        image_paths.append((output_image_path, page_num))
    doc.close()
    return image_paths


def extract_text_and_notes_from_ppt(ppt_path):
    """Extract text and notes from a PowerPoint file."""
    prs = Presentation(ppt_path)
    text_and_notes = []
    for slide in prs.slides:
        slide_text = " ".join(
            [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        )
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ""
        except:
            notes = ""
        text_and_notes.append((slide_text, notes))
    return text_and_notes


async def process_files(
    files: List[UploadFile], research_id: str, reference_id: str
) -> List[Document]:
    """Load and process multiple file types."""
    documents = []
    vectorstore_dir = os.path.join(os.getcwd(), "vectorstore")
    try:
        for file in files:

            file_extension = os.path.splitext(file.filename.lower())[1]
            if file_extension in (".png", ".jpg", ".jpeg"):
                image_content = await file.read()  # Await the read coroutine
                image_text = describe_image(image_content)
                doc = Document(
                    text=f"This is an image({file.filename}) with the following description: \n{image_text}",
                    metadata={
                        "source": file.filename,
                        "type": "image",
                        "caption": image_text,
                        "research_id": research_id,
                        "reference_id": reference_id,
                        "filename": file.filename,
                        "created_at": datetime.utcnow().isoformat(),
                    },
                )
                documents.append(doc)
            elif file_extension == ".pdf":
                try:
                    pdf_documents = await get_pdf_documents(
                        pdf_file=file,
                        research_id=research_id,
                        reference_id=reference_id,
                    )
                    documents.extend(pdf_documents)
                except Exception as e:
                    logger.error(f"Error processing PDF {file.filename}: {e}")
            elif file_extension in (".ppt", ".pptx"):
                try:
                    ppt_documents = process_ppt_file(
                        ppt_path=save_uploaded_file(
                            uploaded_file=file, directory_name="ppt_references"
                        ),
                        research_id=research_id,
                        reference_id=reference_id,
                        filename=file.filename,
                    )
                    documents.extend(ppt_documents)
                except Exception as e:
                    logger.error(f"Error processing PPT {file.filename}: {e}")
            else:
                try:
                    with tempfile.NamedTemporaryFile(
                        delete=True, suffix=f"{file_extension}"
                    ) as temp_file:
                        shutil.copyfileobj(file.file, temp_file)
                        temp_file_path = temp_file.name
                        reader = SimpleDirectoryReader(input_files=[temp_file_path])
                        docs = await reader.aload_data(show_progress=True)
                        for doc in docs:
                            if not hasattr(doc, "metadata"):
                                doc.metadata = {}

                            doc.metadata.update(
                                {
                                    "source": file.filename,
                                    "research_id": research_id,
                                    "reference_id": reference_id,
                                    "filename": file.filename,
                                    "created_at": datetime.utcnow().isoformat(),
                                }
                            )

                        documents.extend(docs)
                except Exception as e:
                    logger.error(
                        f"Error processing {file.filename} using SimpleDirectoryReader: {e}"
                    )

    finally:
        # Delete the vectorstore directory after processing
        if os.path.exists(vectorstore_dir):
            shutil.rmtree(vectorstore_dir)
            logger.info(f"Deleted temporary directory: {vectorstore_dir}")
    return documents
